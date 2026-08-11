"""List slide titles of the current final_presentation_v2 deck."""
from pathlib import Path

from pptx import Presentation

ROOT = Path(__file__).resolve().parents[1]
prs = Presentation(ROOT / "presentation/final_presentation_v3.pptx")
lines = []
for i, slide in enumerate(prs.slides, 1):
    title = ""
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text_frame.text.strip():
            title = shape.text_frame.text.strip().splitlines()[0]
            break
    lines.append(f"{i:2d}. {title}")
(ROOT / "outputs" / "deck_titles.txt").write_text("\n".join(lines), encoding="utf-8")
print("wrote outputs/deck_titles.txt")
